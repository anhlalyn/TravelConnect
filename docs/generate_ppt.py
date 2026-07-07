import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx is not installed. Installing it now...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    COLOR_DARK_BG = RGBColor(15, 23, 42)      # Deep slate/navy
    COLOR_LIGHT_BG = RGBColor(248, 250, 252)  # Light grey
    COLOR_PRIMARY_TEXT = RGBColor(30, 41, 59) # Slate 800
    COLOR_SECONDARY_TEXT = RGBColor(71, 85, 105) # Slate 600
    COLOR_ACCENT = RGBColor(14, 165, 233)     # Light blue/sky
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_CARD_BG = RGBColor(255, 255, 255)
    COLOR_BORDER = RGBColor(226, 232, 240)

    # Image Paths
    IMG_LOGO = "docs/images/logo.png"
    IMG_ERD = "docs/images/erd.png"
    IMG_HOMEPAGE = "docs/images/screenshot_homepage.png"
    IMG_DASHBOARD = "docs/images/screenshot_dashboard.png"

    # Blank layout for custom designs
    blank_slide_layout = prs.slide_layouts[6]

    # Helper: Set background color
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper: Try adding picture safely
    def try_add_picture(slide, img_path, left, top, width=None, height=None):
        if os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, left, top, width=width, height=height)
            except Exception as e:
                print(f"Warning: Could not add picture {img_path}: {e}")
        else:
            print(f"Warning: Picture not found at {img_path}")

    # Helper: Add header to light slides
    def add_slide_header(slide, title_text):
        # Header background shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_DARK_BG
        shape.line.fill.background() # No border
        
        # Title text
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10.0), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.LEFT

        # Small logo in header
        try_add_picture(slide, IMG_LOGO, Inches(12.0), Inches(0.2), width=Inches(0.8), height=Inches(0.8))

    # Helper: Add cards/columns to slides
    def add_card(slide, left, top, width, height, title, items, is_accented=False):
        # Card Background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        if is_accented:
            card.fill.fore_color.rgb = COLOR_DARK_BG
            text_color = COLOR_WHITE
            title_color = COLOR_ACCENT
        else:
            card.fill.fore_color.rgb = COLOR_CARD_BG
            text_color = COLOR_PRIMARY_TEXT
            title_color = COLOR_DARK_BG
            card.line.color.rgb = COLOR_BORDER

        # Card Text Frame
        txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True

        # Card Title
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = title_color
        p_title.font.name = "Arial"
        p_title.space_after = Pt(14)

        # Card Items
        for item in items:
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(14)
            p.font.color.rgb = text_color
            p.font.name = "Arial"
            p.space_after = Pt(8)

    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Dark)
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1, COLOR_DARK_BG)

    # Logo on Title Slide
    try_add_picture(slide1, IMG_LOGO, Inches(1.0), Inches(0.8), width=Inches(1.2), height=Inches(1.2))

    # Subtitle or Topic Indicator
    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "BÁO CÁO CHUYÊN ĐỀ PHÁT TRIỂN HỆ THỐNG"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    p.font.name = "Arial"

    # Main Title
    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(11.333), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "TRAVELCONNECT - NỀN TẢNG MẠNG XÃ HỘI & ĐẶT VÉ DU LỊCH TÍCH HỢP"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = "Arial"

    # Student Info Card
    infoBox = slide1.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(8.0), Inches(2.0))
    tf = infoBox.text_frame
    
    p = tf.paragraphs[0]
    p.text = "Sinh viên thực hiện: "
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_WHITE
    run = p.add_run()
    run.text = "Phan Định Luyện"
    run.font.bold = True
    run.font.color.rgb = COLOR_ACCENT
    p.space_after = Pt(8)

    p2 = tf.add_paragraph()
    p2.text = "Mã số sinh viên: "
    p2.font.size = Pt(16)
    p2.font.color.rgb = COLOR_WHITE
    run2 = p2.add_run()
    run2.text = "22050036"
    run2.font.bold = True
    p2.space_after = Pt(8)

    # Decorator Line
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.8), Inches(3.0), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    # ----------------------------------------------------
    # SLIDE 2: Đặt vấn đề (Problem)
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2, COLOR_LIGHT_BG)
    add_slide_header(slide2, "1. Đặt Vấn Đề & Lý Do Chọn Đề Tài")

    add_card(
        slide2,
        Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.5),
        "Thực trạng & Khó khăn",
        [
            "Nhu cầu du lịch tự túc và tìm kiếm thông tin review chân thực tăng cao.",
            "Sự phân mảnh: Người dùng phải chuyển đổi quá nhiều ứng dụng (App review, App nhắn tin, App đặt vé riêng biệt).",
            "Mất nhiều thời gian để quản lý lịch trình và xác thực giao dịch vé thủ công."
        ]
    )

    add_card(
        slide2,
        Inches(7.133), Inches(2.0), Inches(5.2), Inches(4.5),
        "Giải pháp từ TravelConnect",
        [
            "Mô hình 'All-in-One' tích hợp Mạng xã hội + Đặt vé + Nhắn tin gọi điện realtime.",
            "Tự động hóa hoàn toàn quy trình mua vé, xuất vé QR và check-in nhanh bằng camera.",
            "Giao diện thân thiện kết nối trực tiếp Khách du lịch - Đối tác khu du lịch - Admin."
        ],
        is_accented=True
    )

    # ----------------------------------------------------
    # SLIDE 3: Mục tiêu & Phạm vi (Goal & Scope)
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3, COLOR_LIGHT_BG)
    add_slide_header(slide3, "2. Mục Tiêu & Phạm Vi Đề Tài")

    add_card(
        slide3,
        Inches(1.0), Inches(2.0), Inches(3.5), Inches(4.5),
        "Khách Du Lịch (User)",
        [
            "Tìm kiếm địa điểm du lịch.",
            "Đăng bài viết chia sẻ, tương tác (like, comment, save).",
            "Đặt vé dịch vụ trực tuyến và nhận mã QR check-in.",
            "Gọi điện video, nhắn tin thời gian thực."
        ]
    )

    add_card(
        slide3,
        Inches(4.916), Inches(2.0), Inches(3.5), Inches(4.5),
        "Khu Du Lịch (Business)",
        [
            "Đăng ký thông tin, thiết lập hồ sơ khu du lịch.",
            "Đăng bán các gói vé/dịch vụ.",
            "Quét mã QR của khách trực tiếp bằng camera điện thoại để check-in.",
            "Quản lý booking và xem biểu đồ thống kê doanh thu."
        ]
    )

    add_card(
        slide3,
        Inches(8.833), Inches(2.0), Inches(3.5), Inches(4.5),
        "Hệ Thống (Admin)",
        [
            "Quản lý danh sách người dùng và doanh nghiệp đối tác.",
            "Duyệt các yêu cầu đăng ký kinh doanh.",
            "Quản lý bài viết và tương tác của người dùng.",
            "Thống kê tổng quan hoạt động và cấu hình hệ thống."
        ]
    )

    # ----------------------------------------------------
    # SLIDE 4: Kiến trúc công nghệ (Tech Stack)
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4, COLOR_LIGHT_BG)
    add_slide_header(slide4, "3. Kiến Trúc & Công Nghệ Sử Dụng")

    add_card(
        slide4,
        Inches(1.0), Inches(2.0), Inches(3.5), Inches(4.5),
        "Frontend Client",
        [
            "React.js & Vite: Tốc độ build nhanh, tối ưu hóa giao diện Single Page.",
            "Tailwind CSS: Xây dựng giao diện Responsive, hiện đại và tùy biến cao.",
            "Lucide React: Bộ biểu tượng thiết kế tối giản, trực quan."
        ]
    )

    add_card(
        slide4,
        Inches(4.916), Inches(2.0), Inches(3.5), Inches(4.5),
        "Backend & Database",
        [
            "Node.js & Express.js: Xử lý API tốc độ cao, khả năng chịu tải tốt.",
            "MySQL 8: Cơ sở dữ liệu quan hệ chặt chẽ, tối ưu truy vấn.",
            "JWT & Bcrypt: Bảo mật tài khoản và xác thực API phân quyền."
        ]
    )

    add_card(
        slide4,
        Inches(8.833), Inches(2.0), Inches(3.5), Inches(4.5),
        "Realtime & Deployment",
        [
            "Socket.IO: Truyền tin nhắn thời gian thực và quản lý phòng chat.",
            "WebRTC: Kết nối cuộc gọi thoại, gọi video và livestream peer-to-peer.",
            "Docker & Docker Compose: Đồng bộ môi trường chạy dev & production."
        ],
        is_accented=True
    )

    # ----------------------------------------------------
    # SLIDE 5: Tính năng chính (Features)
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5, COLOR_LIGHT_BG)
    add_slide_header(slide5, "4. Các Phân Hệ Tính Năng Chính")

    add_card(
        slide5,
        Inches(1.0), Inches(2.0), Inches(5.2), Inches(2.1),
        "Xác thực & Mạng xã hội",
        [
            "Đăng ký, gửi mã xác thực OTP thực qua Gmail.",
            "Đăng bài viết kèm ảnh/video, like, comment, lưu trữ."
        ]
    )

    add_card(
        slide5,
        Inches(1.0), Inches(4.4), Inches(5.2), Inches(2.1),
        "Thương mại & Check-in QR",
        [
            "Đặt vé trực tuyến, xuất vé dạng QR Code bảo mật.",
            "Khu du lịch quét QR Code check-in nhanh."
        ]
    )

    add_card(
        slide5,
        Inches(7.133), Inches(2.0), Inches(5.2), Inches(4.5),
        "Realtime Communication & Stream",
        [
            "Nhắn tin tức thời (realtime chat) cá nhân và nhóm.",
            "Đàm thoại trực tiếp (Video Call) công nghệ WebRTC kết nối peer-to-peer mượt mà.",
            "Tích hợp phát Livestream trực tiếp trên nền tảng mạng xã hội kết nối người xem.",
            "Tự động bật/tắt thiết bị camera và micro trực quan."
        ],
        is_accented=True
    )

    # ----------------------------------------------------
    # SLIDE 6: Thiết kế Database (Database Design with ERD)
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6, COLOR_LIGHT_BG)
    add_slide_header(slide6, "5. Thiết Kế Cơ Sở Dữ Liệu & ERD")

    # Left: Database description card
    add_card(
        slide6,
        Inches(1.0), Inches(2.0), Inches(4.8), Inches(4.5),
        "Thực thể chính & Quan hệ",
        [
            "nguoi_dung: Quản lý tài khoản khách hàng, khu du lịch, admin.",
            "bai_viet, binh_luan: Lưu thông tin mạng xã hội.",
            "dat_ve, ve_dich_vu: Quản lý booking vé.",
            "tin_nhan, phong_chat: Nhắn tin realtime.",
            "Ràng buộc khóa ngoại đảm bảo tính toàn vẹn.",
            "Tối ưu hóa các truy vấn JOIN lấy dữ liệu."
        ]
    )

    # Right: ERD image card background
    erd_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.3), Inches(2.0), Inches(6.033), Inches(4.5))
    erd_card.fill.solid()
    erd_card.fill.fore_color.rgb = COLOR_CARD_BG
    erd_card.line.color.rgb = COLOR_BORDER

    # Load and insert ERD diagram
    try_add_picture(slide6, IMG_ERD, Inches(6.4), Inches(2.1), width=Inches(5.833), height=Inches(4.3))

    # ----------------------------------------------------
    # SLIDE 7: Điểm sáng kỹ thuật (Technical Highlights)
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide7, COLOR_LIGHT_BG)
    add_slide_header(slide7, "6. Điểm Sáng Kỹ Thuật & Tối Ưu Hóa")

    add_card(
        slide7,
        Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.5),
        "Hàng đợi ICE Candidate (WebRTC)",
        [
            "Vấn đề: Ứng dụng WebRTC bị lỗi đen màn hình do các gói tin địa chỉ mạng (ICE Candidate) của người gọi đến trước khi người nhận bấm chấp nhận cuộc gọi.",
            "Giải pháp: Xây dựng cơ chế hàng đợi (Queue) lưu trữ tạm thời các ICE candidate.",
            "Kết quả: Ngay sau khi bấm OK kết nối, các gói tin được lấy ra và add đồng loạt, thiết lập camera và âm thanh mượt mà 100%."
        ],
        is_accented=True
    )

    add_card(
        slide7,
        Inches(7.133), Inches(2.0), Inches(5.2), Inches(4.5),
        "Tính ổn định & Triển khai",
        [
            "Gửi mail OTP thực tế: Sử dụng Nodemailer kết nối App Password bảo mật an toàn.",
            "Docker hóa toàn bộ: Dễ dàng cấu hình và khởi chạy ứng dụng với một lệnh duy nhất.",
            "Cơ chế tự khởi tạo Admin mặc định khi database trống, tăng tính tiện lợi khi cài đặt hệ thống."
        ]
    )

    # ----------------------------------------------------
    # SLIDE 8: Demo Giao Diện (UI Mockups with Real screenshots)
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide8, COLOR_LIGHT_BG)
    add_slide_header(slide8, "7. Demo Giao Diện Thực Tế")

    # Left: Homepage Screenshot
    card_left = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.7), Inches(5.2), Inches(5.0))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = COLOR_CARD_BG
    card_left.line.color.rgb = COLOR_BORDER

    # Insert Homepage screenshot
    try_add_picture(slide8, IMG_HOMEPAGE, Inches(1.15), Inches(1.85), width=Inches(4.9), height=Inches(3.3))

    # Homepage Caption
    txBox_left = slide8.shapes.add_textbox(Inches(1.15), Inches(5.2), Inches(4.9), Inches(1.3))
    tf_left = txBox_left.text_frame
    tf_left.word_wrap = True
    p_hl = tf_left.paragraphs[0]
    p_hl.text = "Trang Chủ Mạng Xã Hội & Đặt Vé"
    p_hl.font.size = Pt(16)
    p_hl.font.bold = True
    p_hl.font.color.rgb = COLOR_DARK_BG
    p_hl.font.name = "Arial"
    p_hl.space_after = Pt(4)

    p_hl_desc = tf_left.add_paragraph()
    p_hl_desc.text = "Giao diện chính hiển thị danh sách các bài đăng review du lịch và các địa điểm nổi bật giúp người dùng dễ dàng khám phá và đặt vé trực tuyến."
    p_hl_desc.font.size = Pt(12)
    p_hl_desc.font.color.rgb = COLOR_SECONDARY_TEXT
    p_hl_desc.font.name = "Arial"

    # Right: Dashboard Screenshot
    card_right = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.133), Inches(1.7), Inches(5.2), Inches(5.0))
    card_right.fill.solid()
    card_right.fill.fore_color.rgb = COLOR_CARD_BG
    card_right.line.color.rgb = COLOR_BORDER

    # Insert Dashboard screenshot
    try_add_picture(slide8, IMG_DASHBOARD, Inches(7.283), Inches(1.85), width=Inches(4.9), height=Inches(3.3))

    # Dashboard Caption
    txBox_right = slide8.shapes.add_textbox(Inches(7.283), Inches(5.2), Inches(4.9), Inches(1.3))
    tf_right = txBox_right.text_frame
    tf_right.word_wrap = True
    p_hr = tf_right.paragraphs[0]
    p_hr.text = "Trang Quản Trị Hệ Thống (Dashboard)"
    p_hr.font.size = Pt(16)
    p_hr.font.bold = True
    p_hr.font.color.rgb = COLOR_DARK_BG
    p_hr.font.name = "Arial"
    p_hr.space_after = Pt(4)

    p_hr_desc = tf_right.add_paragraph()
    p_hr_desc.text = "Trang thống kê số liệu tổng quan, doanh thu, quản lý booking, thông tin các khu du lịch đối tác dành cho ban quản trị và đối tác doanh nghiệp."
    p_hr_desc.font.size = Pt(12)
    p_hr_desc.font.color.rgb = COLOR_SECONDARY_TEXT
    p_hr_desc.font.name = "Arial"

    # ----------------------------------------------------
    # SLIDE 9: Đánh giá & Hướng phát triển (Evaluation)
    # ----------------------------------------------------
    slide9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide9, COLOR_LIGHT_BG)
    add_slide_header(slide9, "8. Kết Quả & Hướng Phát Triển")

    add_card(
        slide9,
        Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.5),
        "Đánh giá kết quả",
        [
            "Hoàn thiện 100% các tính năng nghiệp vụ cốt lõi theo kế hoạch đặt ra.",
            "Kết nối chat tức thời và gọi video, livestream hoạt động ổn định và mượt mà.",
            "Hệ thống database được tối ưu truy vấn tốt, đảm bảo tốc độ phản hồi nhanh."
        ]
    )

    add_card(
        slide9,
        Inches(7.133), Inches(2.0), Inches(5.2), Inches(4.5),
        "Hướng phát triển đề tài",
        [
            "Tích hợp các cổng thanh toán điện tử chính thức (VNPay, MoMo, PayPal).",
            "Sử dụng công nghệ AI đề xuất địa điểm du lịch (Recommendation System) theo sở thích.",
            "Phát triển phiên bản ứng dụng di động Hybrid (React Native/Flutter) để tối ưu trải nghiệm scan QR."
        ],
        is_accented=True
    )

    # ----------------------------------------------------
    # SLIDE 10: Conclusion (Dark)
    # ----------------------------------------------------
    slide10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide10, COLOR_DARK_BG)

    # Main Title
    txBox = slide10.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "CẢM ƠN HỘI ĐỒNG THẦY CÔ VÀ CÁC BẠN ĐÃ LẮNG NGHE!"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide10.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(11.333), Inches(1.0))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Em xin kính mong nhận được các ý kiến đóng góp và câu hỏi từ Hội đồng chuyên môn."
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLOR_ACCENT
    p2.font.name = "Arial"
    p2.alignment = PP_ALIGN.CENTER

    # Save presentation
    output_path = "docs/TravelConnect_Slides.pptx"
    try:
        prs.save(output_path)
        print(f"Presentation saved to {output_path} successfully!")
    except PermissionError:
        alternative_path = "docs/TravelConnect_Slides_v2.pptx"
        try:
            prs.save(alternative_path)
            print(f"Permission denied on {output_path} (it is probably open in PowerPoint). Saved to {alternative_path} successfully instead!")
        except Exception as e:
            print(f"Failed to save alternative file: {e}")
    except Exception as e:
        print(f"Error saving presentation: {e}")


if __name__ == "__main__":
    create_presentation()
